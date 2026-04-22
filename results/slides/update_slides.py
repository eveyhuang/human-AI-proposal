#!/usr/bin/env python3
"""Update style_controlled_comparison_presentation.pptx with corrected notebook data."""

from pptx import Presentation

PPTX_PATH = '/Users/eveyhuang/Documents/NICO/human-AI-proposal/results/slides/style_controlled_comparison_presentation.pptx'

prs = Presentation(PPTX_PATH)
slides = prs.slides


def get_shape(slide, shape_id):
    for s in slide.shapes:
        if s.shape_id == shape_id:
            return s
    return None


def rep(shape, *pairs):
    """Replace (old, new) pairs in all runs of a shape."""
    if shape is None or not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            for old, new in pairs:
                if old in run.text:
                    run.text = run.text.replace(old, new)


def set_para(tf, idx, text):
    """Set paragraph idx text to `text`, clearing extra runs."""
    if idx >= len(tf.paragraphs):
        return
    para = tf.paragraphs[idx]
    if not para.runs:
        return
    para.runs[0].text = text
    for r in para.runs[1:]:
        r.text = ''


# ── SLIDE 2 (index 1): Motivation ────────────────────────────────────────────
s = slides[1]
rep(get_shape(s, 20),
    ('Rephrased style AUROC = 0.772', 'Rephrased style AUROC = 0.684'),
    ('Novelty gap (baseline: significant) becomes non-significant after rephrasing',
     'Raw novelty gap persists but attenuates (δ=−0.294, p<0.001); local-density normalized gap disappears (p=0.22)'),
    ('Diversity pattern reverses: humans become more diverse than pooled AI',
     'Diversity pattern reverses: humans become more diverse than pooled AI; Claude-Opus-4.5 shows semantic mode collapse'),
)

# ── SLIDE 7 (index 6): Overview ──────────────────────────────────────────────
s = slides[6]
rep(get_shape(s, 33),
    ('k=10 NN distance to 600 PubMed abstracts', 'k=10 NN distance to 1,030 PubMed articles'),
)

# ── SLIDE 8 (index 7): Key Findings ──────────────────────────────────────────
s = slides[7]

rep(get_shape(s, 13),
    ('δ=−0.885', 'δ=−0.768'),
    ('GPT-5.2 remains homogeneous; Claude and Gemini are intermediate',
     'Claude-Opus-4.5 is maximally homogeneous (mode collapse); GPT-5.2 and Gemini are intermediate'),
)

rep(get_shape(s, 18),
    ('Gap disappears: Novelty difference is no longer significant after rephrasing (δ=−0.132, p=0.17). Style was driving much of the baseline novelty gap.',
     'Partially attenuated: Raw novelty gap persists (δ=−0.294, p<0.001) but is smaller than baseline (δ=−0.34). Local-density normalized gap disappears (δ=+0.016, p=0.22) — humans target sparser literature regions.'),
)

rep(get_shape(s, 23),
    ('NMI=0.389', 'NMI=0.089'),
    ('p<0.001', 'p=0.003'),
    ('Style-reduced AUROC = 0.772', 'Style-reduced AUROC = 0.684'),
    ('Yes, but reduced: Significant cluster segregation remains (NMI=0.389, p<0.001). Style-reduced AUROC = 0.772 vs. 1.000 baseline.',
     'Yes, but weaker: Significant cluster segregation remains (NMI=0.089, p=0.003) but is lower than baseline (0.197), showing style amplified apparent separation. Style-reduced AUROC = 0.684 vs. 1.000 baseline.'),
)

rep(get_shape(s, 28),
    ('GPT (4.13)', 'GPT (4.32)'),
    ('Claude (3.96)', 'Claude (4.01)'),
    ('human (3.39)', 'human (3.59)'),
    ('δ=+0.45', 'δ=+0.933'),
    ('Cross-evaluator reruns confirm GPT advantage.',
     'Cross-evaluator reruns confirm GPT and Claude advantage; Gemini is not significant. Gemini strongly self-deprecates (δ=−0.732).'),
)

rep(get_shape(s, 31),
    ('Human proposals are semantically more varied; AI proposals cluster in different conceptual regions. Quality scoring patterns are robust to style normalization.',
     'Human proposals are semantically more varied; Claude-Opus-4.5 shows near-zero diversity (mode collapse). Raw novelty gap persists. Evaluator biases diverge: GPT self-inflates (δ=+0.933) while Gemini self-deprecates (δ=−0.732).'),
)

# ── SLIDE 9 (index 8): Pairwise Diversity ────────────────────────────────────
s = slides[8]

rep(get_shape(s, 15),
    ('Human mean pairwise distance = 0.364', 'Human mean pairwise distance = 0.443'),
    ('All AI combined = 0.114', 'All AI combined = 0.183'),
    ('δ = −0.885', 'δ = −0.768'),
    ('[−0.323, −0.183]', '[−0.303, −0.211]'),
)

rep(get_shape(s, 18),
    ('GPT-5.2 remains maximally homogeneous', 'Claude-Opus-4.5 is maximally homogeneous (mode collapse)'),
)

rep(get_shape(s, 19),
    ('Mean = 0.030', 'Mean = 0.034'),
    ('Std = 0.006', 'Std = 0.007'),
    ('All GPT proposals occupy essentially the same point',
     'All Claude proposals occupy essentially the same point'),
)

rep(get_shape(s, 22),
    ('Claude and Gemini intermediate', 'GPT-5.2 and Gemini intermediate'),
)

rep(get_shape(s, 23),
    ('Claude mean = 0.144; Gemini mean = 0.160. Both significantly less diverse than humans. Model-specific differences persist after style removal.',
     'GPT-5.2 mean = 0.315 (Std = 0.159); Gemini mean = 0.151 (Std = 0.173). Both significantly less diverse than humans. Claude (mean = 0.034) shows the sharpest collapse.'),
)

# "Key difference from baseline" block — search by content since shape id may vary
for shp in s.shapes:
    if shp.has_text_frame:
        txt = ' '.join(''.join(r.text for r in p.runs) for p in shp.text_frame.paragraphs)
        if 'Key difference from baseline' in txt:
            rep(shp,
                ("Baseline showed Claude > Human in pairwise diversity (Cliff's δ = +1.00) and All AI > Human. After rephrasing, the pattern inverts: Human > All AI. Style was inflating Claude's apparent diversity.",
                 "Baseline showed Claude > Human in pairwise diversity (Cliff's δ = +1.00) and All AI > Human. After rephrasing: Human > GPT-5.2 > Gemini >> Claude (near-zero). Style was entirely masking Claude's semantic homogeneity — Claude collapses to mode collapse after rephrasing."),
            )

# ── SLIDE 10 (index 9): Centroid & NN ────────────────────────────────────────
s = slides[9]

tf31 = get_shape(s, 31).text_frame
set_para(tf31, 0, 'Human proposals are the most dispersed from their centroid (Mean = 0.241), consistent with pairwise results.')
set_para(tf31, 1, 'All AI models are significantly less dispersed than humans (large effects, p < 0.001).')
set_para(tf31, 2, 'Claude-Opus-4.5 is essentially a single point (Mean = 0.016, Std = 0.004) — consistent mode collapse across pairwise, centroid, and NN analyses.')

tf33 = get_shape(s, 33).text_frame
set_para(tf33, 0, 'Human proposals have highest mean NN distance (0.074) and highest outlier rate (30.4%; 7 of 23 proposals).')
set_para(tf33, 1, '87.0% of AI proposals have their nearest neighbor from within the AI group — semantic source-segregation persists after rephrasing.')
set_para(tf33, 2, 'Claude-Opus-4.5 = 0 outliers; 82.6% of its NNs are other Claude proposals — highest same-model NN rate of any group.')

rep(get_shape(s, 35),
    ('Human > Claude/Gemini >> GPT-5.2', 'Human > GPT-5.2 > Gemini >> Claude-Opus-4.5'),
    ('AI same-group NN rate (94.2%)', 'AI same-group NN rate (87.0%)'),
)

# ── SLIDE 11 (index 10): NN Group Membership ─────────────────────────────────
s = slides[10]

rep(get_shape(s, 5),  ('21.7% outliers', '30.4% outliers'))
rep(get_shape(s, 6),  ('65.2% NN from AI', '69.6% NN from AI'))
rep(get_shape(s, 9),  ('8.7% outliers', '0% outliers'))        # Claude
rep(get_shape(s, 10), ('52.2% NN other AI', '82.6% NN same model'))
rep(get_shape(s, 30), ('8.7% outliers', '4.3% outliers'))      # Gemini
rep(get_shape(s, 31), ('60.9% NN other AI', '39.1% NN other AI'))
rep(get_shape(s, 34), ('0% outliers', '8.7% outliers'))        # GPT-5.2
rep(get_shape(s, 35), ('91.3% NN own model', '47.8% NN same model'))

# ── SLIDE 12 (index 11): Embedding Space Visualization ───────────────────────
s = slides[11]

rep(get_shape(s, 12),
    ('5 of 9 total outliers', '7 of 10 total outliers'),
    ('21.7% outlier rate', '30.4% outlier rate'),
)

rep(get_shape(s, 16),
    ('Moderate spread; 2 outlier proposals. After rephrasing, dispersion is reduced compared to baseline — confirming that baseline Claude diversity was partly stylistic.',
     'Near-zero spread; 0 outlier proposals. Claude collapses to a tight cluster after rephrasing — baseline Claude diversity was entirely stylistic (semantic mode collapse).'),
)

rep(get_shape(s, 20),
    ('2 outlier proposals at periphery', '1 outlier proposal at periphery'),
)

rep(get_shape(s, 24),
    ('Tight compact micro-cluster. 0 outliers. Near-zero internal spread — consistent with mode collapse in pairwise and centroid analyses.',
     'Intermediate spread; 2 outlier proposals. GPT-5.2 shows moderate within-group variance (mean pairwise dist = 0.315), distinct from Claude mode collapse.'),
)

# ── SLIDE 13 (index 12): Novelty ─────────────────────────────────────────────
s = slides[12]

rep(get_shape(s, 9),
    ('600 PubMed abstracts', '1,030 PubMed articles'),
)

rep(get_shape(s, 15),
    ('Novelty Gap Disappears After Rephrasing',
     'Raw Novelty Gap Persists but Attenuates; Local-Density Gap Disappears'),
)

rep(get_shape(s, 16),
    ('All AI vs. Human: δ = −0.132 (negligible), permutation p = 0.171. Not significant. Baseline showed δ = −0.34, p = 0.014 (significant human advantage). The gap was largely stylistic.',
     'Raw novelty (AI vs. Human): δ = −0.294 (small), perm p < 0.001 — significant human advantage persists after rephrasing. Local-density normalized: δ = +0.016 (negligible), p = 0.224 — gap disappears when controlling for regional literature density differences. Baseline raw δ = −0.34.'),
)

rep(get_shape(s, 18), ('−0.13', '+0.016'))   # local-density δ stat box

# Label para 0 and p-value para 1 in shape 19
shp19 = get_shape(s, 19)
if shp19 and shp19.has_text_frame:
    rep(shp19,
        ("Cliff's δ (All AI vs Human)", "Local-density δ (All AI vs Human)"),
        ('negligible, p=0.17', 'negligible (local-density), p=0.22'),
    )

rep(get_shape(s, 21), ('−0.005', '−0.030'))   # raw mean diff stat box

shp22 = get_shape(s, 22)
if shp22 and shp22.has_text_frame:
    rep(shp22,
        ('Mean diff (AI − Human)', 'Raw mean diff (AI − Human)'),
        ('[−0.013, +0.003]', '95% CI [−0.051, −0.010]'),
    )

rep(get_shape(s, 26),
    ('"Impact of Heat Wave Stress on Cytokinetic Bottlenecks" — novelty score 0.156. Still highest after rephrasing, suggesting genuine conceptual novelty.',
     '"MaiTool — LLM-powered bioinformatics tools for microbiome analysis" — novelty score 0.240. Followed by "Searching the crosslinking mass spectrometry universe" (0.230).'),
    ('Impact of Heat Wave Stress on Cytokinetic Bottlenecks', 'MaiTool'),
    ('novelty score 0.156', 'novelty score 0.240'),
)

rep(get_shape(s, 30),
    ('All groups have z-scores below the within-literature baseline (z < 0)',
     'All groups have local-density z > 0 (proposals farther from literature than local neighbors). Human z = 1.82 is highest; Gemini = 1.93, GPT = 1.47, Claude = 1.32'),
)

shp35 = get_shape(s, 35)
if shp35 and shp35.has_text_frame:
    rep(shp35,
        ('4 of 15 proposals flagged as both NN-outliers (isolated among proposals) and literature-novel (top-10% by k-NN distance) are Human. Most novel: "Impact of Heat Wave Stress on Cytokinetic Bottlenecks" (z=+0.72).',
         '6 of 10 most literature-novel proposals are Human; 2 Gemini, 1 GPT-5.2. Top NN-outliers also skew human: 7 of 10 total are human. Most novel by raw score: Human "MaiTool" (score=0.240).'),
        ('Impact of Heat Wave Stress on Cytokinetic Bottlenecks', 'MaiTool'),
    )

# ── SLIDE 14 (index 13): Literature Corpus ───────────────────────────────────
s = slides[13]

rep(get_shape(s, 27),
    ('Added topics from topic modeling, 350-> 600', 'Expanded: 350 → 1,030 articles via 35 targeted queries'),
)

shp28 = get_shape(s, 28)
if shp28 and shp28.has_text_frame:
    set_para(shp28.text_frame, 0, '1,030 PubMed articles, 35 targeted queries')
    set_para(shp28.text_frame, 2, 'Within-corpus mean NN dist: 0.110 ± 0.058')

# ── SLIDE 15 (index 14): Literature Space ────────────────────────────────────
s = slides[14]

rep(get_shape(s, 7),  ('600 PubMed articles', '1,030 PubMed articles'))
rep(get_shape(s, 15), ('GPT-5.2 compact cluster', 'Claude-Opus-4.5 compact cluster'))
rep(get_shape(s, 16),
    ('GPT proposals remain tightly clustered in one region of literature space, consistent with the mode collapse seen in diversity analyses. They occupy well-covered literature territory.',
     'Claude proposals remain tightly clustered in one region of literature space, consistent with mode collapse in diversity analyses. GPT-5.2 shows more spread across literature regions.'),
)

# ── SLIDE 16 (index 15): Thematic & Cluster ──────────────────────────────────
s = slides[15]

rep(get_shape(s, 15),
    ('Topic 2 is Human-dominated (95.7% human); Topic 4 is AI-exclusive (0% human). Subsample validation (1,000 resamples): Topics 2 and 4',
     'Topic_2 is Human-dominated (100% of humans participate; Fisher p<0.001); Topic_3 is AI-dominated (4.3% human, Fisher q=0.008). Subsample validation (1,000 resamples): Topic_2 robust (100%), Topic_3 moderate (73%)'),
)

rep(get_shape(s, 17), ('Cluster 0 (n=19)', 'Cluster 0 (n=24)'))
rep(get_shape(s, 18), ('Mixed: 15.8% Human, 84.2% AI', 'Mixed (AI-leaning): 20.8% Human, 79.2% AI'))
rep(get_shape(s, 20), ('Cluster 1 (n=58)', 'Cluster 1 (n=19)'))
rep(get_shape(s, 21), ('AI-dominated: 8.6% Human, 91.4% AI', 'Mixed (Human-leaning): 57.9% Human, 42.1% AI'))
rep(get_shape(s, 23), ('Cluster 2 (n=15)', 'Cluster 2 (n=49)'))
rep(get_shape(s, 24), ('Human-heavy: 60% Human, 40% AI', 'AI-dominated: 14.3% Human, 85.7% AI'))

rep(get_shape(s, 27), ('NMI = 0.389, p < 0.0001', 'NMI = 0.089, p = 0.003'))
rep(get_shape(s, 28), ('ARI = 0.392, p < 0.0001', 'ARI = 0.125, p = 0.001'))
rep(get_shape(s, 29), ('Between/within distance ratio = 1.146, p = 0.001', 'Between/within distance ratio = 1.241, p = 0.002'))

rep(get_shape(s, 32), ('Stronger segregation than baseline', 'Weaker but still significant segregation vs. baseline'))
rep(get_shape(s, 33),
    ('Baseline NMI = 0.197; rephrased NMI = 0.389. After removing style, remaining embedding differences are more cleanly semantic — clusters are better defined. ARI also improves from −0.035 to +0.392.',
     'Baseline NMI = 0.197; rephrased NMI = 0.089. Style was amplifying thematic separation — cluster segregation weakens after rephrasing. ARI improves from −0.035 to +0.125 and the between/within ratio reaches 1.241 (p=0.002), confirming semantic separation persists but is less extreme than baseline.'),
)

# ── SLIDE 17 (index 16): Style Baseline ──────────────────────────────────────
s = slides[16]

rep(get_shape(s, 14), ('0.772', '0.684'))

rep(get_shape(s, 19),
    ('0.772 (±0.096)', '0.684 (±0.102)'),
    ('Permutation p = 0.002', 'Permutation p = 0.023'),
)

rep(get_shape(s, 23),
    ('0.492 ± 0.099', '0.506 ± 0.096'),
    ('0.772 is well above', '0.684 is well above'),
)

rep(get_shape(s, 27),
    ('Type-token ratio (coef = +1.20), avg sentence length (+0.69), and word count (+0.60) are strongest AI predictors. Dash frequency (−0.77) and hedging (−0.30) predict Human.',
     'Type-token ratio (coef = +0.795) is the strongest AI predictor. Number of sentences (coef = −0.536) and hedging rate (−0.369) predict Human. Comma rate (+0.415) and stopword rate (+0.375) also push toward AI.'),
)

rep(get_shape(s, 31), ('AI = 16.6w, Human = 15.7w', 'AI = 18.3w, Human = 18.4w'))
rep(get_shape(s, 33), ('AI = 17.0, Human = 16.5', 'AI ≈ Human (template enforces similar readability)'))
rep(get_shape(s, 35), ('AI = 0.290, Human = 0.292', 'AI = 0.296, Human = 0.293'))

# ── SLIDE 18 (index 17): Style-Adjusted ──────────────────────────────────────
s = slides[17]

rep(get_shape(s, 12),
    ('0.097 (raw)', '0.105 (raw)'),
    ('0.033 (residualized)', '0.064 (residualized)'),
    ('AI same-group NN = 94.2%', 'AI same-group NN = 78.3% in residual space'),
)

rep(get_shape(s, 16),
    ('AI indicator coef = −0.014', 'AI indicator coef = −0.151'),
    ('permutation p = 0.0003', 'permutation p = 0.0002'),
    ('MW δ = −0.3', 'MW δ = −0.764'),
)

rep(get_shape(s, 20),
    ('Outlier threshold (90th pct): 0.741', 'Outlier threshold (90th pct): 0.521'),
    ('Human = 5 outliers (21.7%)', 'Human = 2 outliers (8.7%)'),
    ('Gemini = 1 (4.3%)', 'Gemini = 4 (17.4%)'),
    ('GPT-5.2 gains outliers in residual space while Claude loses them',
     'Gemini gains outliers in residual space (1→4) while Human loses them (7→2)'),
    ("GPT is conceptually isolated, Claude's isolation was stylistic",
     "Gemini is conceptually isolated; Human's peripheral isolation was partly stylistic"),
)

rep(get_shape(s, 23), ('21.7%', '8.7%'))   # Human % in stat bar
rep(get_shape(s, 29), ('4.3%', '17.4%'))   # Gemini % in stat bar

rep(get_shape(s, 37),
    ('94.2% of AI proposals have their nearest neighbor from within the AI group',
     '78.3% of AI proposals have their nearest neighbor from within the AI group'),
)

# ── SLIDE 19 (index 18): Metric–Score Correlation ────────────────────────────
s = slides[18]

rep(get_shape(s, 14),
    ('Relevance (r = −0.585, Holm-corrected), Novelty & Significance (r = −0.498), and Open Science (r = −0',
     'Relevance (centroid_dist r = −0.548) and Novelty & Significance (nn_dist r = −0.498). Main-idea NN distance shows strongest correlations: Relevance r = −0.554, Novelty r = −0.485. Open Science shows weak correlation (r = −0.079)'),
)

# ── SLIDE 21 (index 20): Quality R2 ──────────────────────────────────────────
s = slides[20]

rep(get_shape(s, 13),
    ('GPT-5.2 and Gemini score substantially higher than humans',
     'GPT-5.2 and Claude score substantially higher than humans'),
)

rep(get_shape(s, 14),
    ('GPT-5.2 mean = 4.13 vs. Human 3.39 (δ = −0.72, q < 0.001). Gemini = 3.94 (δ = −0.52, q < 0.001). Claude = 3.96 (not significant). Bootstrap CI (human − GPT): −0.74 [−0.91, −0.57].',
     'GPT-5.2 mean = 4.32 vs. Human 3.59 (δ = −0.72, q < 0.001). Claude = 4.01 (δ = −0.65, q < 0.001). Gemini = 3.83 (not significant vs. human, cross-eval q = 0.68). Bootstrap CI (human − GPT, cross-eval): −0.888 [−1.069, −0.706].'),
)

rep(get_shape(s, 18),
    ('Kruskal-Wallis H = 122.46, p < 2.6×10', 'Kruskal-Wallis H = 98.02, p < 5.2×10'),
    ('10⁻²⁷', '10⁻²²'),
    ('Gemini evaluator: mean = 4.32 (most lenient)', 'Gemini evaluator: mean = 4.30 (most lenient)'),
    ('GPT evaluator: mean = 3.67 (strictest)', 'GPT evaluator: mean = 3.72; Claude = 3.78'),
    ('0.66-point spread', '0.58-point spread (Gemini−GPT)'),
)

rep(get_shape(s, 22),
    ('GPT advantage remains (q < 0.001, δ = −0.87)', 'GPT advantage remains (q < 0.001, δ = −0.992)'),
    ('Gemini attenuates but persists (q = 0.032)', 'Gemini is NOT significant (q = 0.682)'),
    ('Claude not significant', 'Claude advantage confirmed (q < 0.001, δ = −0.652)'),
    ('Human Y1 (3.31)', 'Human Y1 (3.44)'),
    ('Y2 (3.48)', 'Y2 (3.74)'),
)

# ── SLIDE 22 (index 21): Self-Preference R3 ──────────────────────────────────
s = slides[21]

rep(get_shape(s, 13), ('+0.45', '+0.93'))   # GPT delta display

rep(get_shape(s, 14),
    ('Self-inflation (δ=+0.45, q=0.005). Self: 3.88, Other: 3.75.',
     'Strong self-inflation (δ=+0.933, q=3.3×10⁻¹²). Self: 4.00, Other: 3.81.'),
)

rep(get_shape(s, 17), ('~0', '−0.73'))      # Gemini delta display

rep(get_shape(s, 18),
    ('No significant self-preference (q=0.108).',
     'Strong self-DEPRECATION (δ=−0.732, q=7.7×10⁻⁷). Gemini rates its own proposals significantly lower than others\' proposals.'),
)

rep(get_shape(s, 21), ('−0.18', '~0'))      # Claude delta display

rep(get_shape(s, 22),
    ('Slight self-deprecation (q=0.208, ns overall).',
     'Negligible self-preference (δ=−0.053, q=0.717, not significant).'),
)

# Shape 25 / 26: "GPT self-inflates, especially on Relevance and Novelty"
for shp in s.shapes:
    if shp.has_text_frame:
        txt = ' '.join(''.join(r.text for r in p.runs) for p in shp.text_frame.paragraphs)
        if 'Criterion-level tests' in txt:
            rep(shp,
                ('Fixed-effects regression confirms: net self-preference = +0.33 points after controlling for proposal quality and evaluator severity',
                 'Fixed-effects regression confirms net GPT self-preference = +0.33 points. Gemini shows criterion-level self-deprecation on Relevance and Novelty (δ < −0.5, p < 0.001)'),
            )

rep(get_shape(s, 33),
    ('7 Y1 proposals matched', '12 proposals matched'),
    ('p=0.21', 'p=0.26'),
    ('δ=+0.58, p=0.016', 'δ=+0.514, p=0.035'),
)

# ── SLIDE 23 (index 22): Conclusions ─────────────────────────────────────────
s = slides[22]

rep(get_shape(s, 12),
    ('δ=−0.885', 'δ=−0.768'),
    ('Baseline Claude diversity advantage was stylistic; GPT-5.2 mode collapse is semantic.',
     'Baseline Claude diversity advantage was stylistic — Claude-Opus-4.5 shows semantic mode collapse after rephrasing. GPT-5.2 and Gemini are intermediate.'),
)

rep(get_shape(s, 16),
    ('Baseline human novelty advantage disappears after rephrasing (ns, δ=−0.13). The novelty gap was largely a stylistic artifact, not a genuine conceptual difference.',
     'Raw novelty gap persists after rephrasing (δ=−0.294, p<0.001) but is attenuated from baseline. Local-density normalized gap disappears (δ=+0.016, p=0.22). Style was partly responsible; humans also target sparser literature regions.'),
)

rep(get_shape(s, 20),
    ('Semantic segregation is stronger after rephrasing (NMI: 0.197→0.389). Conceptual territory differs between human and AI proposals; style was masking, not creating, this signal.',
     'Semantic segregation weakens after rephrasing (NMI: 0.197→0.089) but remains significant (p=0.003). Style was amplifying apparent thematic separation. No human-dominated cluster; two mixed clusters and one AI-dominated cluster.'),
)

rep(get_shape(s, 24),
    ('GPT-5.2 scores higher than humans even under cross-evaluator control. Self-preference confirmed for GPT. AI reviews remain a valid but homogeneous proxy for human expert feedback.',
     'GPT-5.2 and Claude score higher than humans under cross-evaluator control; Gemini is not significant. GPT strongly self-inflates (δ=+0.933); Gemini strongly self-deprecates (δ=−0.732). AI reviews are valid but evaluator biases must be controlled.'),
)

rep(get_shape(s, 39),
    ('Human proposals explore more varied conceptual territory; AI proposals cluster in specific regions. GPT-5.2 consistently achieves higher AI-assigned quality scores but shows semantic mode collapse — optimizing scores without matching human conceptual breadth.',
     'Human proposals explore more varied conceptual territory; Claude-Opus-4.5 shows near-zero semantic diversity (mode collapse) after style removal. GPT and Claude both achieve higher AI-assigned scores under cross-evaluator control. Large, divergent evaluator biases (GPT inflates, Gemini deflates) must be accounted for in any AI-review pipeline.'),
)

# ── Save ──────────────────────────────────────────────────────────────────────
prs.save(PPTX_PATH)
print("Done. Presentation saved.")
